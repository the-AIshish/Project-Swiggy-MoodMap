from pptx import Presentation
from pptx.util import Pt
import re

def parse_md(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.read().splitlines()
    title = None
    slides = []
    cur = None
    for ln in lines:
        if ln.startswith('# ' ) and title is None:
            title = ln[2:].strip()
            continue
        if ln.startswith('## '):
            if cur:
                slides.append(cur)
            cur = {'title': ln[3:].strip(), 'bullets': []}
            continue
        if cur is not None:
            s = ln.strip()
            if s.startswith('- '):
                cur['bullets'].append(s[2:].strip())
            elif s:
                # treat as paragraph bullet
                cur['bullets'].append(s)
    if cur:
        slides.append(cur)
    return title or 'Swiggy Vibe', slides

prs = Presentation()
md_file = 'Swiggy_Vibe_Stakeholder_Presentation.md'
if __name__ == '__main__':
    title, slides = parse_md(md_file)
    # title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    # Add slides
    for s in slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = s['title']
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(s['bullets']):
            p = body.add_paragraph() if i>0 else body.paragraphs[0]
            p.text = b
            p.font.size = Pt(18)
    out = 'Swiggy_Vibe_Stakeholder_Presentation.pptx'
    prs.save(out)
    print('Wrote', out)
